"""Lightweight file parser: turn any common format into text for LLM analysis.

Deliberately lightweight — mirrors `data_loader`'s philosophy. **No docling,
no torch.** Each format backend is imported lazily, so importing this module
and parsing `.txt`/`.csv` costs nothing; `.pdf`/`.docx`/`.pptx` only pull their
small library when actually parsed. A missing lib raises `ParseError` with an
install hint instead of crashing the caller.

`parse_file(path) -> ParsedDoc{text, kind, meta}`. `kind` is:
- `"document"` — prose (txt/pdf/docx/pptx) → mine for risk points
- `"table"`    — structured rows (csv/xlsx) → business-data insights
`analyze()` (see docanalysis.py) picks its prompt by `kind`.
"""

from __future__ import annotations

import csv
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, List

# Cap how many table rows we feed downstream; huge spreadsheets get summarized
# (the rest are still counted in meta.rows_total). Overridable via DOC_MAX_ROWS,
# bounded by the model's context window (~1500 rows for qwen-plus).
_DEFAULT_MAX_ROWS = 1500


def _max_table_rows() -> int:
    """Row cap, read from settings (env DOC_MAX_ROWS) so it's tunable without code."""
    try:
        from config.settings import settings  # lazy: keep this module config-free at import
        return int(getattr(settings, "doc_max_rows", _DEFAULT_MAX_ROWS) or _DEFAULT_MAX_ROWS)
    except Exception:  # pragma: no cover
        return _DEFAULT_MAX_ROWS


class ParseError(RuntimeError):
    """Unsupported format, or a required parsing library isn't installed."""


@dataclass
class ParsedDoc:
    text: str
    kind: str  # "document" | "table"
    meta: Dict[str, Any] = field(default_factory=dict)


def parse_file(path) -> ParsedDoc:
    """Parse a file by extension into a ParsedDoc. Raises ParseError on
    unsupported formats or missing optional libs."""
    p = Path(path)
    ext = p.suffix.lower()
    if ext in (".txt", ".md"):
        return _read_text(p)
    if ext == ".csv":
        return _read_csv(p)
    if ext in (".xlsx", ".xlsm"):
        return _read_excel(p)
    if ext == ".pdf":
        return _read_pdf(p)
    if ext == ".docx":
        return _read_docx(p)
    if ext == ".doc":
        return _read_doc(p)
    if ext == ".pptx":
        return _read_pptx(p)
    raise ParseError(
        f"不支持的文件格式: {ext or '(无扩展名)'}"
        "（支持 txt/md/csv/xlsx/pdf/docx/pptx）"
    )


# --- text ------------------------------------------------------------------

def _read_text(p: Path) -> ParsedDoc:
    text = p.read_text(encoding="utf-8", errors="replace")
    return ParsedDoc(
        text=text, kind="document",
        meta={"format": "text", "filename": p.name, "chars": len(text)},
    )


# --- tables (csv / xlsx) ---------------------------------------------------

def _read_csv(p: Path) -> ParsedDoc:
    # Windows/Excel 导出的中文 CSV 多为 GBK/GB18030，纯英文多为 UTF-8（带或不带
    # BOM）。逐个尝试，latin-1 兜底（永不抛错，最坏是乱码但不崩）。
    import io
    raw = p.read_bytes()
    text = None
    for enc in ("utf-8-sig", "gb18030", "utf-8", "latin-1"):
        try:
            text = raw.decode(enc)
            break
        except UnicodeDecodeError:
            continue
    rows = list(csv.reader(io.StringIO(text)))
    return _rows_to_table_doc(p, rows, fmt="csv")


def _read_excel(p: Path) -> ParsedDoc:
    try:
        from openpyxl import load_workbook  # lazy, like data_loader._read_excel
    except ImportError as exc:  # pragma: no cover
        raise ParseError(f"解析 {p.name} 需要 openpyxl：pip install openpyxl") from exc
    wb = load_workbook(p, data_only=True)
    ws = wb.active
    sheet_title = ws.title
    rows = [[("" if c is None else str(c)) for c in row]
            for row in ws.iter_rows(values_only=True)]
    wb.close()
    doc = _rows_to_table_doc(p, rows, fmt="xlsx", extra={"sheet": sheet_title})
    return doc


def _rows_to_table_doc(p: Path, rows: List[List[str]], *, fmt: str,
                       extra: Dict[str, Any] = None) -> ParsedDoc:
    width = max((len(r) for r in rows), default=0) or 1
    norm = [(list(r) + [""] * width)[:width] for r in rows]
    header = norm[0] if norm else []
    body = norm[1:]
    total = len(body)
    cap = _max_table_rows()
    shown = body[:cap]

    lines: List[str] = []
    if header:
        lines.append("| " + " | ".join(header) + " |")
        lines.append("|" + "|".join("---" for _ in header) + "|")
    for r in shown:
        lines.append("| " + " | ".join(r) + " |")
    text = "\n".join(lines)

    meta: Dict[str, Any] = {
        "format": fmt,
        "filename": p.name,
        "columns": [h for h in header if h],
        "rows_total": total,
        "rows_shown": len(shown),
        "truncated": total > len(shown),
    }
    if extra:
        meta.update(extra)
    return ParsedDoc(text=text, kind="table", meta=meta)


# --- documents (pdf / docx / pptx) ----------------------------------------

def _read_pdf(p: Path) -> ParsedDoc:
    try:
        import pdfplumber  # lazy
    except ImportError as exc:
        raise ParseError(
            f"解析 {p.name} 需要 pdfplumber：pip install pdfplumber"
            "（注：扫描件纯图片 PDF 需 OCR，本期暂不支持）"
        ) from exc
    pages: List[str] = []
    n_pages = 0
    with pdfplumber.open(p) as pdf:
        n_pages = len(pdf.pages)
        for pg in pdf.pages:
            t = pg.extract_text() or ""
            if t.strip():
                pages.append(t.strip())
    text = "\n\n".join(pages)
    return ParsedDoc(
        text=text, kind="document",
        meta={"format": "pdf", "filename": p.name, "pages": n_pages, "chars": len(text)},
    )


def _read_docx(p: Path) -> ParsedDoc:
    try:
        import docx  # lazy (python-docx)
    except ImportError as exc:
        raise ParseError(f"解析 {p.name} 需要 python-docx：pip install python-docx") from exc
    d = docx.Document(str(p))
    chunks: List[str] = [par.text for par in d.paragraphs if par.text and par.text.strip()]
    for t in d.tables:
        for row in t.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                chunks.append(" | ".join(cells))
    text = "\n".join(chunks)
    return ParsedDoc(
        text=text, kind="document",
        meta={"format": "docx", "filename": p.name, "chars": len(text)},
    )


def _read_doc(p: Path) -> ParsedDoc:
    # 旧版二进制 .doc（OLE2 复合文档）。用 olefile 打开、读 WordDocument 流提取正文。
    # 无完美轻量解析器；中文 .doc 正文多为 UTF-16-LE，西文多为 cp1252，两种都试。
    try:
        import olefile  # lazy (optional dep)
    except ImportError as exc:
        raise ParseError(f"解析 {p.name} 需要 olefile：pip install olefile") from exc
    try:
        ole = olefile.OleFileIO(str(p))
        try:
            data = ole.openstream("WordDocument").read()
        finally:
            ole.close()
    except Exception as exc:
        raise ParseError(
            f"{p.name} 不是有效的 .doc（可能实为 RTF/HTML 改名或已损坏）。"
            "建议在 Word 里另存为 .docx 后重传。"
        ) from exc
    text = _extract_doc_text(data)
    if len(text.strip()) < 8:
        raise ParseError(
            f"{p.name} 未能提取到正文（格式特殊或损坏）。建议另存为 .docx 或 PDF 后重传。"
        )
    return ParsedDoc(
        text=text, kind="document",
        meta={"format": "doc", "filename": p.name, "chars": len(text)},
    )


def _extract_doc_text(data: bytes) -> str:
    """Crude text extraction from a .doc WordDocument stream.

    Decodes the stream as both UTF-16-LE and cp1252, keeps runs of printable
    chars (incl. CJK), and returns the longer result. Noisy but good enough
    for LLM analysis; the model tolerates surrounding formatting noise.
    """
    import re
    run_re = re.compile("[" + chr(0x20) + "-" + chr(0x9FFF) + "]{4,}")

    def _runs(s: str) -> str:
        return "\n".join(run_re.findall(s))

    cand_utf16 = _runs(data.decode("utf-16-le", errors="ignore"))
    cand_cp = _runs(data.decode("cp1252", errors="ignore"))
    # collapse blank lines / excessive whitespace
    best = cand_utf16 if len(cand_utf16) >= len(cand_cp) else cand_cp
    return re.sub(r"\n{3,}", "\n\n", best).strip()


def _read_pptx(p: Path) -> ParsedDoc:
    try:
        from pptx import Presentation  # lazy (python-pptx, optional)
    except ImportError as exc:
        raise ParseError(f"解析 {p.name} 需要 python-pptx：pip install python-pptx") from exc
    prs = Presentation(str(p))
    slides_text: List[str] = []
    for i, slide in enumerate(prs.slides, 1):
        bits: List[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs)
                    if t.strip():
                        bits.append(t.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        bits.append(" | ".join(cells))
        if bits:
            slides_text.append(f"[Slide {i}]\n" + "\n".join(bits))
    text = "\n\n".join(slides_text)
    return ParsedDoc(
        text=text, kind="document",
        meta={"format": "pptx", "filename": p.name, "slides": len(prs.slides), "chars": len(text)},
    )


__all__ = ["parse_file", "ParsedDoc", "ParseError"]
